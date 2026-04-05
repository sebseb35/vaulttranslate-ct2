from io import BytesIO
from pathlib import Path

import pytest
from pptx import Presentation

from packages.adapter_pptx import PptxDocumentAdapter
from packages.core.models import DocumentFormat, Segment


FIXTURES_DIR = Path(__file__).resolve().parents[1] / "fixtures" / "pptx"


def _fixture_bytes(name: str) -> bytes:
    return (FIXTURES_DIR / name).read_bytes()


def test_pptx_adapter_extracts_shapes_runs_tables() -> None:
    adapter = PptxDocumentAdapter()
    source = _fixture_bytes("sample.pptx")

    segments = adapter.extract_segments(source)
    texts = [segment.text for segment in segments]

    assert adapter.supported_format is DocumentFormat.PPTX
    assert all(segment.segment_id.startswith("pptx-") for segment in segments)
    assert "Hello " in texts
    assert "world" in texts
    assert "Second paragraph" in texts
    assert "Cell A1" in texts
    assert "Cell B2" in texts


def test_pptx_adapter_rebuild_roundtrip_preserves_text_runs_and_tables() -> None:
    adapter = PptxDocumentAdapter()
    source = _fixture_bytes("sample.pptx")
    extracted = adapter.extract_segments(source)

    translated = tuple(
        Segment(segment_id=segment.segment_id, text=f"{segment.text} FR")
        for segment in extracted
    )
    rebuilt_bytes = adapter.rebuild_document(source, translated)

    rebuilt = Presentation(BytesIO(rebuilt_bytes))
    slide = rebuilt.slides[0]

    title_runs = slide.shapes[0].text_frame.paragraphs[0].runs
    assert title_runs[0].text == "Hello  FR"
    assert title_runs[1].text == "world FR"
    assert title_runs[1].font.bold is True
    assert slide.shapes[1].text_frame.paragraphs[0].text == "Second paragraph FR"
    assert slide.shapes[2].table.cell(0, 0).text == "Cell A1 FR"
    assert slide.shapes[2].table.cell(1, 1).text == "Cell B2 FR"


def test_pptx_adapter_rebuild_validates_segment_count_and_ids() -> None:
    adapter = PptxDocumentAdapter()
    source = _fixture_bytes("sample.pptx")
    extracted = adapter.extract_segments(source)

    with pytest.raises(ValueError, match="same number of segments"):
        adapter.rebuild_document(source, extracted[:-1])

    wrong_first = (Segment(segment_id="pptx-9999", text="bad"),) + extracted[1:]
    with pytest.raises(ValueError, match="segment mismatch"):
        adapter.rebuild_document(source, wrong_first)
