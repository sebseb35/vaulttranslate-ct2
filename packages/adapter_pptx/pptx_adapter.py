from __future__ import annotations

from collections.abc import Sequence
from dataclasses import dataclass
from io import BytesIO
from typing import Any, cast

from pptx import Presentation
from pptx.presentation import Presentation as PptxPresentation
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.table import _Cell
from pptx.text.text import _Paragraph, _Run

from packages.core import DocumentAdapter, DocumentFormat, Segment


@dataclass(frozen=True, slots=True)
class _TextTarget:
    kind: str
    node: _Run | _Paragraph


class PptxDocumentAdapter(DocumentAdapter):
    @property
    def supported_format(self) -> DocumentFormat:
        return DocumentFormat.PPTX

    def extract_segments(self, document_bytes: bytes) -> tuple[Segment, ...]:
        presentation = _load_presentation(document_bytes)
        targets = _collect_text_targets(presentation)
        return tuple(
            Segment(
                segment_id=f"pptx-{index}",
                text=_target_text(target),
            )
            for index, target in enumerate(targets)
        )

    def rebuild_document(
        self, original_document: bytes, translated_segments: Sequence[Segment]
    ) -> bytes:
        presentation = _load_presentation(original_document)
        targets = _collect_text_targets(presentation)

        if len(translated_segments) != len(targets):
            raise ValueError("PPTX rebuild requires the same number of segments as extraction")

        for index, (target, translated_segment) in enumerate(
            zip(targets, translated_segments, strict=True)
        ):
            expected_id = f"pptx-{index}"
            if translated_segment.segment_id != expected_id:
                raise ValueError(
                    "PPTX segment mismatch: "
                    f"expected '{expected_id}' but got '{translated_segment.segment_id}'"
                )
            _set_target_text(target, translated_segment.text)

        output = BytesIO()
        presentation.save(output)
        return output.getvalue()


def _load_presentation(document_bytes: bytes) -> PptxPresentation:
    try:
        return Presentation(BytesIO(document_bytes))
    except Exception as exc:  # pragma: no cover - defensive handling
        raise ValueError("Failed to parse PPTX document bytes") from exc


def _collect_text_targets(presentation: PptxPresentation) -> list[_TextTarget]:
    targets: list[_TextTarget] = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            _collect_shape_text_targets(shape, targets)

    return targets


def _collect_shape_text_targets(shape: BaseShape, targets: list[_TextTarget]) -> None:
    if isinstance(shape, GroupShape):
        for child in shape.shapes:
            _collect_shape_text_targets(child, targets)
        return

    if getattr(shape, "has_table", False):
        table = cast(Any, shape).table
        for row in table.rows:
            for cell in row.cells:
                _collect_cell_text_targets(cell, targets)

    if getattr(shape, "has_text_frame", False):
        text_frame = cast(Any, shape).text_frame
        _collect_paragraph_targets(text_frame.paragraphs, targets)


def _collect_cell_text_targets(cell: _Cell, targets: list[_TextTarget]) -> None:
    _collect_paragraph_targets(cell.text_frame.paragraphs, targets)


def _collect_paragraph_targets(
    paragraphs: Sequence[_Paragraph],
    targets: list[_TextTarget],
) -> None:
    for paragraph in paragraphs:
        runs = [run for run in paragraph.runs if run.text != ""]
        if runs:
            for run in runs:
                targets.append(_TextTarget(kind="run", node=run))
            continue

        if paragraph.text != "":
            targets.append(_TextTarget(kind="paragraph", node=paragraph))


def _target_text(target: _TextTarget) -> str:
    if target.kind == "run":
        return target.node.text
    if target.kind == "paragraph":
        return target.node.text
    raise ValueError(f"Unsupported target kind: {target.kind}")


def _set_target_text(target: _TextTarget, text: str) -> None:
    if target.kind == "run":
        target.node.text = text
        return
    if target.kind == "paragraph":
        target.node.text = text
        return
    raise ValueError(f"Unsupported target kind: {target.kind}")
